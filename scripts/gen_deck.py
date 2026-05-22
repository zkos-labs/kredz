#!/usr/bin/env python3
"""Generate the kredz.xyz pitch deck for the Build on Canton Hackathon."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Colors ──
DARK      = RGBColor(0x00, 0x04, 0x1F)
DARK_HERO = RGBColor(0x00, 0x04, 0x1F)
ACCENT    = RGBColor(0xB5, 0x69, 0x39)
GOLD      = RGBColor(0xD4, 0xA8, 0x53)
LIGHT     = RGBColor(0xEF, 0xF4, 0xFF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x89, 0x8A, 0x92)
DARK_CARD = RGBColor(0x0A, 0x0F, 0x2A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def make_slide(bg_color=DARK_HERO):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    return slide

def add_text(slide, text, left, top, width, height, size=14, color=LIGHT,
             bold=False, font='Calibri', align=PP_ALIGN.LEFT, is_title=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Georgia' if is_title else font
    p.alignment = align
    return txBox

def add_rich_text(slide, segments, left, top, width, height, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """segments: list of (text, size, color, bold, font) tuples"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for i, seg in enumerate(segments):
        if i > 0:
            run = p.add_run()
        else:
            run = p.runs[0] if p.runs else p.add_run()
        text, size, color, bold, font = seg
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font
    return txBox

def add_bullets(slide, items, left, top, width, height, size=14, color=LIGHT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.level = 0
    return txBox

def add_accent_bar(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, bg=DARK_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(0.5)
    return shape

# ═══════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════
s1 = make_slide()

add_text(s1, "KREDZ", 1.0, 1.0, 11, 1.8, size=72, color=GOLD, bold=True, font='Georgia', align=PP_ALIGN.LEFT)
add_accent_bar(s1, 1.0, 2.9, 2.0, 0.06, GOLD)
add_text(s1, "The Dual-Privacy Credit Identity Layer", 1.0, 3.2, 11, 1.0, size=28, color=LIGHT, font='Calibri', bold=False)
add_text(s1, "Bridging crypto-native borrowers to institutional lenders\nwith ZK proofs on Midnight + sub-transaction privacy on Canton", 1.0, 4.2, 11, 1.2, size=16, color=MUTED, font='Calibri')

add_text(s1, "Build on Canton Hackathon · June 2026", 1.0, 6.0, 6, 0.4, size=12, color=ACCENT)
add_text(s1, "Track: Private DeFi & Capital Markets", 1.0, 6.4, 6, 0.4, size=12, color=MUTED)

# ═══════════════════════════════════════════════════════════
# SLIDE 2: PROBLEM
# ═══════════════════════════════════════════════════════════
s2 = make_slide()
add_accent_bar(s2, 1.0, 0.8, 0.08, 0.6, GOLD)
add_text(s2, "The Credit Identity Gap", 1.3, 0.7, 10, 0.8, size=36, color=WHITE, bold=True, font='Georgia', is_title=True)

# Stats cards
for i, (num, label) in enumerate([
    ("1.4B+", "Adults globally unbanked\n— no formal credit file"),
    ("$0", "Under-collateralized\nlending on-chain today"),
    ("3 in 4", "Crypto-native users lack\ninstitutional credit access"),
]):
    x = 1.0 + i * 3.9
    add_card(s2, x, 1.8, 3.5, 2.2)
    add_text(s2, num, x + 0.3, 2.0, 2.9, 0.8, size=40, color=GOLD, bold=True)
    add_text(s2, label, x + 0.3, 2.9, 2.9, 1.0, size=14, color=LIGHT)

add_text(s2, "Existing Solutions Fall Short", 1.0, 4.4, 5, 0.5, size=18, color=ACCENT, bold=True)
add_bullets(s2, [
    "No protocol combines web2 banking + web3 on-chain + literacy into one portable score",
    "All lending requires either full identity disclosure (bad for privacy) or lending blind (bad for compliance)",
    "Institutional lenders entering DeFi need compliance-ready credit oracles — none exist",
], 1.0, 4.9, 11.5, 1.8, size=13, color=LIGHT)

add_text(s2, "The regulatory moment: MiCA (2024) + GENIUS Act (2025) mandate compliance for on-chain lending.", 1.0, 6.8, 11.5, 0.4, size=11, color=MUTED)

# ═══════════════════════════════════════════════════════════
# SLIDE 3: SOLUTION
# ═══════════════════════════════════════════════════════════
s3 = make_slide()
add_accent_bar(s3, 1.0, 0.8, 0.08, 0.6, GOLD)
add_text(s3, "Three-Layer KREDZ Score (0–1000)", 1.3, 0.7, 10, 0.8, size=36, color=WHITE, bold=True, font='Georgia', is_title=True)

layers = [
    ("Layer 1: On-Chain Signals", "40%", ACCENT,
     "Wallet age · DeFi interactions · Repayment history\nTransaction frequency · Governance · Cross-chain activity\n→ Pulled from Midnight indexer + EVM RPCs"),
    ("Layer 2: ZK-Attested KYC", "40%", GOLD,
     "Income range · Employment status · Bank account history\nE-wallet volume · Credit commitments · Jurisdiction\n→ ZK-proven on Midnight — raw data never exposed"),
    ("Layer 3: Financial Literacy", "20%", RGBColor(0x5C, 0xCE, 0x7A),
     "5+ modules with quizzes, XP, and time-decay\nStreak bonuses · First-attempt accuracy · Recency\n→ Behavioral moat — longitudinal, hard to fake"),
]
for i, (title, weight, color, desc) in enumerate(layers):
    y = 1.8 + i * 1.8
    add_card(s3, 1.0, y, 11.3, 1.5, DARK_CARD)
    w = add_text(s3, weight, 11.3, y + 0.15, 0.8, 0.6, size=28, color=color, bold=True, align=PP_ALIGN.RIGHT)
    add_text(s3, title, 1.3, y + 0.1, 5, 0.5, size=18, color=color, bold=True)
    add_text(s3, desc, 1.3, y + 0.6, 9.5, 0.8, size=12, color=LIGHT)

add_accent_bar(s3, 1.0, 7.1, 11.3, 0.02, ACCENT)
add_text(s3, "AI scoring engine: XGBoost + Python · Attested on-chain via Midnight ZK circuits · Ported to Canton via KredzScore DAML", 0.8, 7.15, 11.5, 0.3, size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 4: ARCHITECTURE
# ═══════════════════════════════════════════════════════════
s4 = make_slide()
add_accent_bar(s4, 1.0, 0.8, 0.08, 0.6, GOLD)
add_text(s4, "Three-Network Architecture", 1.3, 0.7, 10, 0.8, size=36, color=WHITE, bold=True, font='Georgia', is_title=True)
add_text(s4, "Midnight builds the score privately. Canton distributes to institutions. Base makes it portable across EVM.", 1.0, 1.5, 11.5, 0.5, size=14, color=MUTED)

networks = [
    ("MIDNIGHT", "Credit Identity", "ZK proofs via Compact circuits\nSelective disclosure\nLiteracy XP on-chain\nShielded credential storage", "Dark"),
    ("CANTON", "Institutional Credit", "Sub-transaction privacy\nKredzScore DAML + QueryScore\nImmutable KredzAuditLog\nLender subscriptions\nMiCA/GENIUS Act compliant", "Light"),
    ("BASE", "EVM Portability", "ERC-8004 SBT badge\nIKredzOracle for DeFi\nECDSA attestation verifier\nCross-chain composability", "Light"),
]

for i, (name, role, detail, style) in enumerate(networks):
    x = 1.0 + i * 4.0
    bg = DARK_CARD if i == 0 else DARK_CARD
    add_card(s4, x, 2.2, 3.6, 3.8, bg)
    header_color = GOLD if i == 0 else (ACCENT if i == 1 else RGBColor(0x5C, 0xCE, 0x7A))
    add_text(s4, name, x + 0.2, 2.4, 3.2, 0.5, size=20, color=header_color, bold=True, font='Georgia')
    add_text(s4, role, x + 0.2, 2.9, 3.2, 0.4, size=12, color=LIGHT, bold=True)
    add_text(s4, detail, x + 0.2, 3.4, 3.2, 2.4, size=12, color=LIGHT)

# arrows between
for i in range(2):
    ax = 4.6 + i * 4.0
    add_text(s4, "→", ax, 3.5, 0.8, 0.5, size=28, color=GOLD, align=PP_ALIGN.CENTER)

add_text(s4, "→ Users onboard once on Midnight. Their score flows to wherever lenders operate.", 1.0, 6.4, 11.5, 0.5, size=12, color=GOLD, align=PP_ALIGN.CENTER, bold=True)

# ═══════════════════════════════════════════════════════════
# SLIDE 5: DUAL PRIVACY
# ═══════════════════════════════════════════════════════════
s5 = make_slide()
add_accent_bar(s5, 1.0, 0.8, 0.08, 0.6, GOLD)
add_text(s5, "Dual-Privacy Model — Our Hackathon Differentiator", 1.3, 0.7, 10, 0.8, size=36, color=WHITE, bold=True, font='Georgia', is_title=True)

add_card(s5, 1.0, 1.8, 5.3, 5.0)
add_text(s5, "Zero-Knowledge (Midnight)", 1.3, 2.0, 4.7, 0.5, size=20, color=GOLD, bold=True)
add_text(s5, "Borrower Privacy", 1.3, 2.5, 4.7, 0.3, size=14, color=MUTED)
add_bullets(s5, [
    "Prove financial attributes\nwithout revealing raw data",
    "Income is \"above $5,000\"\nnot \"$7,250\"",
    "Selective disclosure:\nchoose what to prove",
    "Compact circuits on Midnight\nNo proof generation overhead\non Canton",
], 1.3, 3.0, 4.7, 3.0, size=12, color=LIGHT)

add_card(s5, 7.0, 1.8, 5.3, 5.0)
add_text(s5, "Sub-Transaction (Canton)", 7.3, 2.0, 4.7, 0.5, size=20, color=ACCENT, bold=True)
add_text(s5, "Lender Confidentiality", 7.3, 2.5, 4.7, 0.3, size=14, color=MUTED)
add_bullets(s5, [
    "Lenders only see score queries\nthey initiated themselves",
    "No other participant\ncan observe the transaction",
    "DAML signatories:\nonly kredz + querying lender",
    "KredzAuditLog: immutable\ncompliance evidence\nco-signed by both parties",
], 7.3, 3.0, 4.7, 3.0, size=12, color=LIGHT)

add_text(s5, "Two complementary privacy models. One portable score.", 1.0, 7.0, 11.5, 0.4, size=14, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 6: CANTON DEEP DIVE
# ═══════════════════════════════════════════════════════════
s6 = make_slide()
add_accent_bar(s6, 1.0, 0.8, 0.08, 0.6, GOLD)
add_text(s6, "Canton Deep Dive — DAML Contracts", 1.3, 0.7, 10, 0.8, size=36, color=WHITE, bold=True, font='Georgia', is_title=True)

contracts = [
    ("KredzScore", "Score registry with QueryScore choice\nLender queries → KredzScoreResponse\nUpdateScore, ExpireScore choices"),
    ("KredzQuery", "Response + Audit log\nKredzScoreResponse: visible to operator+lender only\nKredzAuditLog: immutable, co-signed, never archived"),
    ("KredzSubscription", "Lender watchlist\nAddBorrower, RemoveBorrower\nWebhook delivery\nConfigurable score threshold"),
]

for i, (name, detail) in enumerate(contracts):
    x = 1.0 + i * 4.0
    add_card(s6, x, 1.8, 3.6, 2.8)
    add_text(s6, name, x + 0.2, 2.0, 3.2, 0.4, size=18, color=GOLD, bold=True, font='Georgia')
    add_text(s6, detail, x + 0.2, 2.5, 3.2, 2.0, size=12, color=LIGHT)

add_text(s6, "Why Canton?", 1.0, 5.0, 5, 0.5, size=18, color=ACCENT, bold=True)
add_bullets(s6, [
    "Sub-transaction privacy: QueryScore results are hidden from \u2003all non-participating parties — no other lender can see who you're assessing",
    "Immutable compliance: KredzAuditLog is co-signed, satisfies MiCA/GENIUS Act audit requirements",
    "Institutional-ready: Goldman, BNP, Deloitte already operate Canton nodes — kredz plugs directly into their infrastructure",
], 1.0, 5.5, 11.5, 1.6, size=13, color=LIGHT)

add_text(s6, "5 DAML Script tests written: lifecycle · update · expiry · subscription · sub-transaction privacy", 1.0, 7.0, 11.5, 0.4, size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 7: DEMO FLOW
# ═══════════════════════════════════════════════════════════
s7 = make_slide()
add_accent_bar(s7, 1.0, 0.8, 0.08, 0.6, GOLD)
add_text(s7, "End-to-End Demo Flow", 1.3, 0.7, 10, 0.8, size=36, color=WHITE, bold=True, font='Georgia', is_title=True)

steps = [
    ("1", "Connect Wallet", "User connects Lace\n(Midnight) + MetaMask\n(Base)"),
    ("2", "Select Tier", "Choose privacy level:\nAnonymous, Pseudonymous,\nor Full Compliance"),
    ("3", "Build Score", "On-chain signals + ZK\nKYC attestation +\nliteracy modules →\nKREDZ Score (0-1000)"),
    ("4", "Attest", "Score published to\nMidnight attestation\nstore via ZK proof"),
    ("5", "Sync to Canton", "Backend detects attestation\n→ creates KredzScore\nDAML contract on Canton"),
    ("6", "Lender Queries", "Institutional lender\nqueries score via\nQueryScore → response\nvisible only to lender!"),
]

for i, (num, title, desc) in enumerate(steps):
    x = 0.5 + i * 2.1
    add_card(s7, x, 1.8, 1.9, 4.5)
    add_text(s7, num, x + 0.15, 2.0, 0.5, 0.5, size=28, color=GOLD, bold=True)
    add_text(s7, title, x + 0.15, 2.6, 1.6, 0.4, size=14, color=LIGHT, bold=True)
    add_text(s7, desc, x + 0.15, 3.1, 1.6, 3.0, size=11, color=LIGHT)
    if i < 5:
        add_text(s7, "→", x + 1.95, 3.5, 0.3, 0.5, size=20, color=ACCENT, bold=True)

add_text(s7, "Privacy verified at every step: ZK on Midnight protects borrower data. Sub-transaction privacy on Canton protects lender activity.", 0.5, 6.6, 12.5, 0.5, size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 8: PRIVACY TIERS
# ═══════════════════════════════════════════════════════════
s8 = make_slide()
add_accent_bar(s8, 1.0, 0.8, 0.08, 0.6, GOLD)
add_text(s8, "Three Privacy Tiers", 1.3, 0.7, 10, 0.8, size=36, color=WHITE, bold=True, font='Georgia', is_title=True)

tiers = [
    ("Tier 0: Anonymous", "Score: 0–400", "No identity required\n\n• On-chain wallet analysis\n• Transaction history\n• DeFi interactions\n• Literacy modules\n\n→ Micro-lending access\n→ No KYC friction"),
    ("Tier 1: Pseudonymous", "Score: 0–650", "1 ZK-proven attribute\n\n• All Tier 0 features\n• Income OR employment\n  ZK attestation\n• Mid-tier lending\n\n→ Privacy preserved\n→ Higher credit limits"),
    ("Tier 2: Full Compliance", "Score: 0–1000", "Full ZK-KYC bundle\n\n• All Tier 1 features\n• Full credential attestation\n• Canton lender access\n\n→ Institutional liquidity\n→ MiCA/GENIUS compliant"),
]

for i, (title, score, desc) in enumerate(tiers):
    x = 1.0 + i * 4.0
    border = GOLD if i == 1 else MUTED
    add_card(s8, x, 1.8, 3.6, 5.0)
    border_color = GOLD if i == 1 else MUTED
    card = add_card(s8, x, 1.8, 3.6, 5.0)
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5 if i == 1 else 0.5)
    add_text(s8, title, x + 0.2, 2.0, 3.2, 0.4, size=18, color=GOLD, bold=True)
    add_text(s8, score, x + 0.2, 2.4, 3.2, 0.3, size=14, color=ACCENT, bold=True)
    add_text(s8, desc, x + 0.2, 2.9, 3.2, 3.8, size=12, color=LIGHT)

# ═══════════════════════════════════════════════════════════
# SLIDE 9: REVENUE
# ═══════════════════════════════════════════════════════════
s9 = make_slide()
add_accent_bar(s9, 1.0, 0.8, 0.08, 0.6, GOLD)
add_text(s9, "Revenue Model — $240K Year 1 Target", 1.3, 0.7, 10, 0.8, size=36, color=WHITE, bold=True, font='Georgia', is_title=True)

streams = [
    ("Score API (B2B)", "$15K/mo", "$90K/yr",
     "Institutional lenders pay per score query.\nTiered: $0.50 (≤1K), $0.25 (≤10K), $0.10 (10K+).\nCanton direct: enterprise negotiated."),
    ("Premium Subscriptions", "$10K/mo", "$60K/yr",
     "Users: $7/mo (coaching, priority KYC, 15 modules)\n$15/mo (adverse action simulator, score forecast).\n→ Built-in literacy gamification drives conversion."),
    ("Literacy SaaS", "$10K/mo", "$60K/yr",
     "White-label literacy modules for neobanks/wallets.\n$5K/mo per institution.\nCustom module creation: $10K/module."),
    ("Protocol Integrations", "$5K/mo", "$30K/yr",
     "Base oracle read: free (open).\nWebhook subscription: $500/mo.\nCustom support: $5K one-time."),
]

for i, (name, mo, yr, desc) in enumerate(streams):
    y = 1.8 + i * 1.3
    add_card(s9, 1.0, y, 11.3, 1.1)
    add_text(s9, name, 1.3, y + 0.05, 3, 0.4, size=16, color=GOLD, bold=True)
    add_text(s9, mo, 4.5, y + 0.15, 1.5, 0.4, size=20, color=ACCENT, bold=True, align=PP_ALIGN.RIGHT)
    add_text(s9, yr, 6.2, y + 0.15, 1.5, 0.4, size=14, color=MUTED, align=PP_ALIGN.RIGHT)
    add_text(s9, desc, 1.3, y + 0.45, 10.5, 0.6, size=11, color=LIGHT)

add_text(s9, "Total Year 1: $240,000 revenue · 500 credentialed users · 5 institutional partners · 10 Base DeFi integrations", 0.5, 7.0, 12.5, 0.4, size=13, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 10: ROADMAP
# ═══════════════════════════════════════════════════════════
s10 = make_slide()
add_accent_bar(s10, 1.0, 0.8, 0.08, 0.6, GOLD)
add_text(s10, "Roadmap — From Testnet to Institutional Scale", 1.3, 0.7, 10, 0.8, size=36, color=WHITE, bold=True, font='Georgia', is_title=True)

phases = [
    ("Phase 0: Now – Jul 2026", "Testnet + Canton Launch",
     "• Compile kredz.compact\n• Canton DevNet deployment\n• Score API live\n• DAML Script tests passing\n• Hackathon submission"),
    ("Phase 1: Q3 2026", "Midnight Mainnet MVP",
     "• Mainnet deployment\n• 100 credentialed users\n• 3 Base DeFi integrations\n• SumSub KYC live\n• Brankas SEA bank data"),
    ("Phase 2: Q4 2026", "Canton Institutional Lenders",
     "• First institutional lender\n  on Canton\n• KredzScore DAML contracts\n  deployed to MainNet\n• Compliance audit export\n• MiCA documentation"),
    ("Phase 3: 2027", "Scale + ZK-ML",
     "• 10+ institutional lenders\n• ZK-ML proof for scoring\n  engine (EZKL)\n• 15 literacy modules\n• Literacy SaaS licensing\n• DAO governance"),
]

for i, (phase, title, desc) in enumerate(phases):
    y = 1.8 + i * 1.35
    add_card(s10, 1.0, y, 11.3, 1.15)
    add_text(s10, phase, 1.3, y + 0.05, 3, 0.3, size=12, color=ACCENT, bold=True)
    add_text(s10, title, 1.3, y + 0.3, 4, 0.3, size=16, color=GOLD, bold=True)
    add_text(s10, desc, 1.3, y + 0.6, 10.5, 0.5, size=11, color=LIGHT)

add_text(s10, "Current: kredz-frontend built · DAML contracts written · Scoring engine scaffolded · Ready for Canton hackathon deployment", 0.5, 7.0, 12.5, 0.4, size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 11: TEAM
# ═══════════════════════════════════════════════════════════
s11 = make_slide()
add_accent_bar(s11, 1.0, 0.8, 0.08, 0.6, GOLD)
add_text(s11, "Team & Technology", 1.3, 0.7, 10, 0.8, size=36, color=WHITE, bold=True, font='Georgia', is_title=True)

add_text(s11, "Builder", 1.0, 1.8, 5, 0.5, size=18, color=ACCENT, bold=True)
add_text(s11, "DeDanzi / Tawf Labs\n@dedanzi/midnight-mobile-sdk (npm)\nMidnight Hackathon participant\nActive in Midnight Discord community", 1.0, 2.3, 5.5, 2.0, size=14, color=LIGHT)

add_text(s11, "Built on Midnight Build Club", 7.0, 1.8, 5, 0.5, size=18, color=GOLD, bold=True)
add_bullets(s11, [
    "Midnight Network — ZK privacy",
    "Canton Network — Institutional DeFi",
    "Base — EVM portability",
    "Compact + DAML + Solidity",
], 7.0, 2.3, 5.5, 2.0, size=13, color=LIGHT)

add_text(s11, "Tech Stack", 1.0, 4.5, 5, 0.5, size=18, color=ACCENT, bold=True)
add_bullets(s11, [
    "React 19 + Vite + TypeScript + Tailwind v4",
    "Framer Motion animations",
    "Node.js + Express + Python XGBoost backend",
    "Foundry (Solidity) + Hardhat",
    "Docker Compose (Canton LocalNet)",
], 1.0, 5.0, 5.5, 2.0, size=13, color=LIGHT)

add_text(s11, "Open Source", 7.0, 4.5, 5, 0.5, size=18, color=GOLD, bold=True)
add_bullets(s11, [
    "github.com/kredz-labs/kredz",
    "DAML contracts in canton/daml/",
    "5 DAML Script tests",
    "OpenAPI 3.1 Score API spec",
    "Docker Compose for Canton LocalNet",
], 7.0, 5.0, 5.5, 2.0, size=13, color=LIGHT)

# ═══════════════════════════════════════════════════════════
# SLIDE 12: CLOSING
# ═══════════════════════════════════════════════════════════
s12 = make_slide()
add_text(s12, "Thank You", 1.0, 1.5, 11, 1.5, size=72, color=GOLD, bold=True, font='Georgia', align=PP_ALIGN.CENTER)
add_accent_bar(s12, 5.5, 3.2, 2.3, 0.06, GOLD)
add_text(s12, "kredz.xyz — The Dual-Privacy Credit Identity Layer", 1.0, 3.6, 11.5, 1.0, size=24, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s12, "One score. Two privacy models. Three networks. Zero raw data exposed.", 1.0, 4.5, 11.5, 0.5, size=16, color=MUTED, align=PP_ALIGN.CENTER)

add_text(s12, "github.com/kredz-labs/kredz", 1.0, 5.8, 11.5, 0.5, size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(s12, "Built for the Build on Canton Hackathon · June 2026", 1.0, 6.5, 11.5, 0.4, size=12, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s12, "Track: Private DeFi & Capital Markets · Canton Foundation", 1.0, 6.9, 11.5, 0.4, size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ── Save ──
output_path = "/home/zidan/Documents/Github/kredz/kredz-pitch-deck.pptx"
prs.save(output_path)
print(f"Pitch deck saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
